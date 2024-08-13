from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import os
from io import BytesIO
from mimetypes import guess_type
from pymongo import MongoClient
from gridfs import GridFS
from motor.motor_asyncio import AsyncIOMotorClient
from bson import ObjectId
from fastapi.responses import JSONResponse, StreamingResponse
import pypandoc
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from fpdf import FPDF
import random
import string
import json
from pymongo.errors import PyMongoError
from pydantic import BaseModel
# from moviepy.editor import VideoFileClip
# from quiz_generator import export_quiz
# from flash_card_generator import export_flashcards
# from summary_generator import export_summary
from Summary.export_summary import export_summary
from Quiz.export_quiz import export_quiz
from FlashCards.export_flashcards import export_flashcards
from gemini import prompt_everyting
from gemini import prompt_recitation_comparison
import speech_recognition as sr
from pydub import AudioSegment
# from speech_to_text import get_audio

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize the MongoDB client
client = MongoClient('mongodb://localhost:27017/')
db = client['MYTUTOR']
fs = GridFS(db)
metadata_collection = db["metadata"]  # Collection for storing metadata



@app.post("/uploadtodb")
async def upload_file(file: UploadFile = File(...), responseData: str = Form(...)):
    try:
        # Read the file contents
        contents = await file.read()
        
        # Parse the responseData from JSON string
        try:
            response_data = json.loads(responseData)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON for response data.")
        
        # Store the file in GridFS
        file_id = fs.put(contents, filename=file.filename)

        # Store the response_data in a separate MongoDB collection
        metadata_collection.insert_one({
            "file_id": str(file_id),
            "response_data": response_data
        })

        # Return the response with file info and response data
        return JSONResponse(content={"file_id": str(file_id), "filename": file.filename, "response_data": response_data})

    except PyMongoError as e:
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File upload failed: {str(e)}")

@app.get("/files/{file_id}")
async def get_file(file_id: str):
    try:
        gridout = await fs.get(ObjectId(file_id))
        headers = {
            'Content-Disposition': f'attachment; filename="{gridout.filename}"'
        }
        return StreamingResponse(gridout, headers=headers)

    except Exception as e:
        raise HTTPException(status_code=404, detail="File not found")

@app.get("/recentfiles")
async def get_recent_files():
    try:
        files_collection = db['fs.files']
        
        # Find the most recent files
        cursor = files_collection.find({}, {"_id": 1, "filename": 1}).sort("_id", -1).limit(10)
        files = [file for file in cursor]

        # Fetch the file metadata
        file_details = []
        for file in files:
            file_id = file['_id']
            file_details.append({
                'id': f"{file_id}",
                'filename': file['filename'],
                'file_url': f"/view/{file_id}"  # URL to view the file
            })

        return JSONResponse(content=file_details, status_code=200)
    except Exception as e:
        print(f"Error: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@app.get("/view/{file_id}")
async def view_file(file_id: str):
    try:
        file_id = ObjectId(file_id)
        grid_out = fs.get(file_id)
        
        # Get the file MIME type
        mime_type, _ = guess_type(grid_out.filename)
        if mime_type is None:
            mime_type = "application/octet-stream"  # Default MIME type
        
        # Stream the file to the client
        return StreamingResponse(BytesIO(grid_out.read()), media_type=mime_type)
    except Exception as e:
        print(f"Error: {str(e)}")
        raise HTTPException(status_code=404, detail="File not found")
    
@app.get("/filemetadata/{file_id}")
async def get_file_metadata(file_id: str):
    try:
        # Ensure that file_id is valid
        if not ObjectId.is_valid(file_id):
            raise HTTPException(status_code=400, detail="Invalid file ID format.")
        
        # Fetch file metadata from the database
        metadata = metadata_collection.find_one({"file_id": file_id})

        if metadata:
            return JSONResponse(content={
                "file_id": metadata.get("file_id"),
                "filename": metadata.get("filename"),
                "file_url": metadata.get("file_url"),
                "response_data": metadata.get("response_data", {})  # Ensure response_data is included
            })
        else:
            raise HTTPException(status_code=404, detail="File metadata not found.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error retrieving file metadata: {str(e)}")


def handle_pdf(file_path):
    reader = PdfReader(file_path)
    number_of_pages = len(reader.pages)
    s = ""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        s += text
    return s

def handle_txt(file_path):
    with open(file_path, "r") as f:
        s = f.read()
    return s

def handle_docx(file_path):
    d = Document(file_path)
    s = ""
    for paragraph in d.paragraphs:
        s += paragraph.text + "\n"
    return s

# def convert_mp3_to_wav(mp3_file, wav_file):
#     audio = AudioSegment.from_mp3(mp3_file)
#     audio.export(wav_file, format="wav")

# def handle_mp3(file_path):
#     wav_file = file_path.replace(".mp3", ".wav")
    
#     # Convert MP3 to WAV
#     convert_mp3_to_wav(file_path, wav_file)

#     recognizer = sr.Recognizer()
    
#     try:
#         with sr.AudioFile(wav_file) as source:
#             audio = recognizer.record(source)
#             try:
#                 # Using Google's Web Speech API (free)
#                 transcript = recognizer.recognize_google(audio)
#                 return transcript
#             except sr.UnknownValueError:
#                 return "Google Speech Recognition could not understand the audio"
#             except sr.RequestError:
#                 return "Could not request results from Google Speech Recognition service"
#     finally:
#         # Clean up WAV file
#         if os.path.exists(wav_file):
#             os.remove(wav_file)


# def handle_mp4(file_path):
#     video = VideoFileClip(file_path)
#     audio = video.audio
#     audio_file_path = "file.mp3"
#     audio.write_audiofile(audio_file_path)
#     audio.close()
#     video.close()
#     return handle_mp3(audio_file_path)

def handle_pptx(file_path):
    p = Presentation(file_path)
    s = ""
    for slide in p.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                s += shape.text + "\n"
    return s

@app.get("/")
async def hello_world():
    return {"message": "Hello, World!"}

class RecitationRequest(BaseModel):
    provided_text: str
    transcript: str

@app.post("/recitation")
async def recitation(request: Request):
    try:
        data = await request.json()
        provided_text = data.get('provided_text')
        transcript = data.get('transcript')

        # Assuming you have a function to compare the texts
        comparison_report = prompt_recitation_comparison(provided_text, transcript)

        return {"comparison_report": comparison_report}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error comparing transcript: {str(e)}")



@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    name = file.filename
    extension = name.split(".")[-1]
    file_path = f"file.{extension}"
    
    with open(file_path, "wb") as buffer:
        buffer.write(file.file.read())

    if extension == "pdf":
        s = handle_pdf(file_path)
    elif extension == "txt":
        s = handle_txt(file_path)
    elif extension == "docx":
        s = handle_docx(file_path)
    # elif extension == "mp3":
    #     s = handle_mp3(file_path)
    # elif extension == "mp4":
    #     s = handle_mp4(file_path)
    elif extension == "pptx":
        s = handle_pptx(file_path)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    response = prompt_everyting(s)
    
    return response

@app.post("/export")
async def export(request: Request):
    req = await request.json()
    selected = req.get("selected")
    data = req.get("data")
    
    filename = ""
    if selected == 0:
        filename = "Summary.docx"
        export_summary(data, filename)
    elif selected == 1:
        filename = "Flashcards.docx"
        export_flashcards(data, filename)
    else:
        filename = "Quiz.docx"
        export_quiz(data, filename)
    
    return FileResponse(path=filename, filename=filename, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=5000)
